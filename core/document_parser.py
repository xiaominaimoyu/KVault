from pathlib import Path
from dataclasses import dataclass


MAX_FILE_SIZE = 512 * 1024 * 1024  # 512 MB


@dataclass
class ParsedDocument:
    content: str
    metadata: dict


class DocumentParser:
    SUPPORTED = {".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx"}

    def parse(self, file_path: str) -> ParsedDocument:
        path = Path(file_path)
        ext = path.suffix.lower()
        if ext not in self.SUPPORTED:
            raise ValueError(f"不支持的文件格式: {ext}，仅支持 {', '.join(sorted(self.SUPPORTED))}")

        size = path.stat().st_size
        if size > MAX_FILE_SIZE:
            raise ValueError(f"文件过大: {size / 1024 / 1024:.1f} MB，超过限制 {MAX_FILE_SIZE / 1024 / 1024:.0f} MB")

        parser = self._dispatch(ext)
        if parser is None:
            raise ValueError(f"不支持的文件格式: {ext}")
        return parser(file_path)

    def supported_formats(self) -> list[str]:
        return sorted(self.SUPPORTED)

    def _dispatch(self, ext: str):
        return {
            ".txt": self._parse_txt,
            ".md": self._parse_txt,
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
        }.get(ext)

    def _parse_txt(self, path: str) -> ParsedDocument:
        text = Path(path).read_text("utf-8", errors="ignore")
        return ParsedDocument(text, {"source": path})

    def _parse_pdf(self, path: str) -> ParsedDocument:
        import fitz
        with fitz.open(path) as doc:
            text = "\n".join(page.get_text() for page in doc)
            return ParsedDocument(text, {"source": path, "pages": doc.page_count})

    def _parse_docx(self, path: str) -> ParsedDocument:
        from docx import Document
        doc = Document(path)
        lines: list[str] = []

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            if p.style and p.style.name and p.style.name.startswith("Heading"):
                try:
                    level = int(p.style.name.replace("Heading", ""))
                    lines.append(f"{'#' * level} {text}")
                except ValueError:
                    lines.append(text)
            else:
                lines.append(text)

        for table in doc.tables:
            lines.append("[Table]")
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    lines.append(row_text)

        return ParsedDocument("\n".join(lines), {"source": path})

    def _parse_xlsx(self, path: str) -> ParsedDocument:
        from openpyxl import load_workbook
        with load_workbook(path, data_only=True) as wb:
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        lines.append(row_text)
        return ParsedDocument("\n".join(lines), {"source": path})

    def _parse_pptx(self, path: str) -> ParsedDocument:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return ParsedDocument("\n\n".join(slides), {"source": path})
